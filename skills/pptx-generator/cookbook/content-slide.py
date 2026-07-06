# /// layout
# purpose = "Title + bullet points — the fallback layout when no visual layout fits"
# best_for = "Process steps; lists of 3-6 items with roughly equal weight; text-heavy content"
# avoid_when = "You have numbers/metrics (use stats-slide); you have 2 things to compare (use two-column-slide); you have 3-5 equal cards (use multi-card-slide). Use this LAST."
# max_title_chars = 70
# max_bullets = 6
# max_bullet_chars = 90
# instructions = """
# Title top-left in brand primary. Bullets left-aligned below.
# Max 6 bullets. Each bullet should be ONE idea — ruthlessly short.
# Sub-bullets only if unavoidable; never go 3 levels deep.
# Content-slide should be <25% of any deck. If you're using it more, rethink the deck.
# """
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_content_slide(prs, title, bullets,
                      bg="#07090F", primary="#00B4D8", text="#E0E0E0", accent="#8ECAE6"):
    """
    bullets: list of str, or list of (str, [sub_bullets]) tuples
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(bg)

    W, H = prs.slide_width, prs.slide_height

    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), W - Inches(1.0), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = hex_to_rgb(primary)

    # Divider line
    from pptx.util import Emu
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.25), W - Inches(0.5), Inches(1.25))
    line.line.color.rgb = hex_to_rgb(accent)
    line.line.width = Emu(12700)  # 1pt

    # Bullets
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), W - Inches(1.0), H - Inches(2.0))
    tf2 = tx2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets[:6]):
        if isinstance(bullet, tuple):
            text_val, subs = bullet
        else:
            text_val, subs = bullet, []

        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p2.level = 0
        r2 = p2.add_run()
        r2.text = f"• {text_val}"
        r2.font.size = Pt(18)
        r2.font.color.rgb = hex_to_rgb(text)
        p2.space_after = Pt(4)

        for sub in subs:
            ps = tf2.add_paragraph()
            ps.level = 1
            rs = ps.add_run()
            rs.text = f"  – {sub}"
            rs.font.size = Pt(14)
            rs.font.color.rgb = hex_to_rgb(accent)

    return slide
