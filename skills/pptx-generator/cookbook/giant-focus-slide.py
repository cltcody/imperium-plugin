# /// layout
# purpose = "One big idea, word, or short phrase that needs to land with maximum force"
# best_for = "Section transitions; single-word impact ('Why now?'); announcing a decision; a number so important it gets its own slide"
# avoid_when = "More than 1-3 words in the hero element; you need supporting context (add a subtitle or use stats-slide)"
# max_hero_chars = 40
# max_subtitle_chars = 100
# instructions = """
# Giant text centred on slide — fills 60-70% of the width.
# Optional single-line subtitle in smaller secondary colour.
# Nothing else. Maximum whitespace. The emptiness amplifies the message.
# """
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_giant_focus_slide(prs, hero_text, subtitle="",
                          bg="#07090F", hero_color="#00B4D8", sub_color="#8ECAE6"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(bg)

    W, H = prs.slide_width, prs.slide_height

    # Auto-size font based on length
    if len(hero_text) <= 6:
        font_size = Pt(80)
    elif len(hero_text) <= 15:
        font_size = Pt(60)
    else:
        font_size = Pt(44)

    v_centre = H / 2 - Inches(0.6) if not subtitle else H / 2 - Inches(0.9)

    # Hero text
    tx = slide.shapes.add_textbox(Inches(0.8), v_centre, W - Inches(1.6), Inches(1.4))
    tf = tx.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = hero_text
    r.font.size = font_size
    r.font.bold = True
    r.font.color.rgb = hex_to_rgb(hero_color)

    # Subtitle
    if subtitle:
        ts = slide.shapes.add_textbox(Inches(1.0), v_centre + Inches(1.5),
                                      W - Inches(2.0), Inches(0.6))
        tfs = ts.text_frame
        ps = tfs.paragraphs[0]
        ps.alignment = PP_ALIGN.CENTER
        rs = ps.add_run()
        rs.text = subtitle
        rs.font.size = Pt(20)
        rs.font.color.rgb = hex_to_rgb(sub_color)

    return slide
