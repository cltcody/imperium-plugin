# /// layout
# purpose = "3-5 equal-weight items shown as individual cards — capabilities, pillars, features, steps"
# best_for = "Product capabilities; solution pillars; process steps; benefit statements"
# avoid_when = "Fewer than 3 items (use content-slide or two-column-slide); items have very different amounts of content; items need hierarchical nesting"
# max_cards = 5
# max_card_title_chars = 30
# max_card_body_chars = 120
# instructions = """
# Horizontal row of 3-5 cards. Each card: icon area (optional) + title + short body.
# Cards should have a subtle border or background — not floating text boxes.
# All cards same height and width. Consistent spacing between them.
# """
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_multi_card_slide(prs, cards, title="",
                         bg="#07090F", primary="#00B4D8",
                         card_bg="#0D1B2A", card_border="#00B4D8",
                         title_color="#E0E0E0", body_color="#8ECAE6"):
    """
    cards: list of {"title": str, "body": str} dicts (3-5)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(bg)

    W, H = prs.slide_width, prs.slide_height
    cards = cards[:5]
    n = len(cards)

    # Optional deck title
    top_offset = Inches(0.35)
    if title:
        tx = slide.shapes.add_textbox(Inches(0.5), top_offset, W - Inches(1.0), Inches(0.6))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = hex_to_rgb(primary)
        top_offset = Inches(1.1)

    margin = Inches(0.4)
    gap = Inches(0.2)
    card_w = (W - 2 * margin - (n - 1) * gap) / n
    card_h = H - top_offset - Inches(0.5)

    for i, card in enumerate(cards):
        cx = margin + i * (card_w + gap)
        cy = top_offset

        # Card background
        rect = slide.shapes.add_shape(1, cx, cy, card_w, card_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = hex_to_rgb(card_bg)
        rect.line.color.rgb = hex_to_rgb(card_border)
        rect.line.width = Emu(12700)

        # Accent bar at top of card
        bar = slide.shapes.add_shape(1, cx, cy, card_w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(primary)
        bar.line.fill.background()

        # Card title
        tt = slide.shapes.add_textbox(cx + Inches(0.12), cy + Inches(0.15),
                                      card_w - Inches(0.24), Inches(0.55))
        tft = tt.text_frame
        tft.word_wrap = True
        pt = tft.paragraphs[0]
        pt.alignment = PP_ALIGN.LEFT
        rt = pt.add_run()
        rt.text = card["title"]
        rt.font.size = Pt(14)
        rt.font.bold = True
        rt.font.color.rgb = hex_to_rgb(title_color)

        # Card body
        tb = slide.shapes.add_textbox(cx + Inches(0.12), cy + Inches(0.78),
                                      card_w - Inches(0.24), card_h - Inches(1.0))
        tfb = tb.text_frame
        tfb.word_wrap = True
        pb = tfb.paragraphs[0]
        pb.alignment = PP_ALIGN.LEFT
        rb = pb.add_run()
        rb.text = card.get("body", "")
        rb.font.size = Pt(12)
        rb.font.color.rgb = hex_to_rgb(body_color)

    return slide
