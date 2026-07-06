# /// layout
# purpose = "A single powerful quote, customer statement, or executive insight that needs maximum impact"
# best_for = "Customer pain quotes from discovery; analyst or exec endorsements; 'before' pain statement to open a section"
# avoid_when = "The quote is longer than 3 sentences; there's no clear attribution; you need to present data alongside it"
# max_quote_chars = 220
# max_attribution_chars = 60
# instructions = """
# Large quote in centre of slide, brand accent colour, quotation marks oversized.
# Attribution (name, title, company) in smaller secondary colour below.
# Nothing else on the slide. Negative space is intentional. Let the quote breathe.
# """
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_quote_slide(prs, quote, attribution="",
                    bg="#07090F", quote_color="#00B4D8",
                    mark_color="#1A3A4A", attr_color="#8ECAE6"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(bg)

    W, H = prs.slide_width, prs.slide_height

    # Opening quotation mark
    tq1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(1.2), Inches(1.2))
    tfq1 = tq1.text_frame
    pq1 = tfq1.paragraphs[0]
    rq1 = pq1.add_run()
    rq1.text = "“"
    rq1.font.size = Pt(96)
    rq1.font.color.rgb = hex_to_rgb(mark_color)

    # Quote text
    top = Inches(1.5) if len(quote) < 100 else Inches(1.3)
    tx = slide.shapes.add_textbox(Inches(0.9), top, W - Inches(1.8), H - top - Inches(1.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = quote
    r.font.size = Pt(22 if len(quote) > 150 else 26)
    r.font.italic = True
    r.font.color.rgb = hex_to_rgb(quote_color)

    # Attribution
    if attribution:
        ta = slide.shapes.add_textbox(Inches(0.9), H - Inches(1.3), W - Inches(1.8), Inches(0.6))
        tfa = ta.text_frame
        pa = tfa.paragraphs[0]
        pa.alignment = PP_ALIGN.CENTER
        ra = pa.add_run()
        ra.text = f"— {attribution}"
        ra.font.size = Pt(14)
        ra.font.color.rgb = hex_to_rgb(attr_color)

    return slide
