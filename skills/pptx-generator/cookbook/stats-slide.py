# /// layout
# purpose = "2-4 large metrics or statistics that need to stand out visually"
# best_for = "ROI numbers; KPI summaries; 'before vs after' metrics; value drivers with $/$% figures"
# avoid_when = "More than 4 stats (split across 2 slides); no actual numbers (use content-slide)"
# max_stats = 4
# max_stat_value_chars = 12
# max_stat_label_chars = 40
# instructions = """
# Large number in brand accent colour, label beneath in secondary.
# Lay out horizontally (2-4 cards equally spaced).
# Optional title at top. Optional thin divider.
# The NUMBER is the hero — make it as large as possible.
# """
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_stats_slide(prs, stats, title="",
                    bg="#07090F", primary="#00B4D8", value_color="#00B4D8",
                    label_color="#8ECAE6", title_color="#E0E0E0"):
    """
    stats: list of {"value": "90%", "label": "Classification accuracy"} dicts (max 4)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(bg)

    W, H = prs.slide_width, prs.slide_height
    stats = stats[:4]
    n = len(stats)

    # Optional title
    if title:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), W - Inches(1.0), Inches(0.6))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = hex_to_rgb(title_color)

    top = Inches(2.0) if title else Inches(1.8)
    card_w = (W - Inches(1.0)) / n
    gap = Inches(0.5) / n

    for i, stat in enumerate(stats):
        x = Inches(0.5) + i * (card_w + gap)

        # Value
        tv = slide.shapes.add_textbox(x, top, card_w - gap, Inches(1.6))
        tfv = tv.text_frame
        pv = tfv.paragraphs[0]
        pv.alignment = PP_ALIGN.CENTER
        rv = pv.add_run()
        rv.text = stat["value"]
        rv.font.size = Pt(54)
        rv.font.bold = True
        rv.font.color.rgb = hex_to_rgb(value_color)

        # Label
        tl = slide.shapes.add_textbox(x, top + Inches(1.7), card_w - gap, Inches(0.8))
        tfl = tl.text_frame
        tfl.word_wrap = True
        pl = tfl.paragraphs[0]
        pl.alignment = PP_ALIGN.CENTER
        rl = pl.add_run()
        rl.text = stat["label"]
        rl.font.size = Pt(14)
        rl.font.color.rgb = hex_to_rgb(label_color)

    return slide
