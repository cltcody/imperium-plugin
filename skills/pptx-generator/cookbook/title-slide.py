# /// layout
# purpose = "Opening slide for a presentation — company / deck name, subtitle, optional date"
# best_for = "First slide of any deck; section dividers with a major heading"
# avoid_when = "Content slides; anything with bullet points or data"
# max_title_chars = 60
# max_subtitle_chars = 120
# instructions = """
# Large centred title in brand primary colour on brand background.
# Subtitle in muted/secondary colour below.
# Optional date or account name in small text at bottom-right.
# No bullet points. No body copy. Title is the whole message.
# """
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_title_slide(prs, title, subtitle="", footer="",
                    bg="#07090F", primary="#00B4D8", secondary="#8ECAE6"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = hex_to_rgb(bg)

    W, H = prs.slide_width, prs.slide_height

    # Title — centred, large
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), W - Inches(1.6), Inches(1.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(primary)

    # Subtitle
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), W - Inches(2.0), Inches(1.0))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(18)
        r2.font.color.rgb = hex_to_rgb(secondary)

    # Footer
    if footer:
        tx3 = slide.shapes.add_textbox(W - Inches(3.5), H - Inches(0.7), Inches(3.0), Inches(0.4))
        tf3 = tx3.text_frame
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.RIGHT
        r3 = p3.add_run()
        r3.text = footer
        r3.font.size = Pt(10)
        r3.font.color.rgb = hex_to_rgb(secondary)

    return slide
