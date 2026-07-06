# /// layout
# purpose = "Compare or contrast two things side by side — before/after, us vs. them, option A vs. option B"
# best_for = "Competitive comparison; current state vs. future state; problem vs. solution"
# avoid_when = "More than 2 columns needed (use multi-card-slide); items aren't comparable; one side has much more content"
# max_title_chars = 70
# max_column_header_chars = 30
# max_bullets_per_column = 5
# instructions = """
# Title at top. Two columns below, each with a header and 3-5 bullet points.
# Use a subtle vertical divider between columns.
# Left column often = current state / competitor / problem.
# Right column often = future state / us / solution.
# Optionally highlight the right column with brand accent border to signal 'the better option'.
# """
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_two_column_slide(prs, title, left_header, left_bullets,
                         right_header, right_bullets,
                         bg="#07090F", primary="#00B4D8", text="#E0E0E0",
                         accent="#8ECAE6", highlight_right=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(bg)

    W, H = prs.slide_width, prs.slide_height

    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), W - Inches(1.0), Inches(0.65))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = hex_to_rgb(primary)

    col_w = (W - Inches(1.4)) / 2
    left_x = Inches(0.5)
    right_x = left_x + col_w + Inches(0.4)
    col_top = Inches(1.2)

    def add_column(x, header, bullets, highlight=False):
        # Header background rect (subtle)
        hdr_rect = slide.shapes.add_shape(1, x, col_top, col_w, Inches(0.55))
        hdr_rect.fill.solid()
        hdr_rect.fill.fore_color.rgb = hex_to_rgb(primary if highlight else "#1E2A35")
        hdr_rect.line.fill.background()

        # Header text
        th = slide.shapes.add_textbox(x + Inches(0.1), col_top + Inches(0.05), col_w - Inches(0.2), Inches(0.45))
        tfh = th.text_frame
        ph = tfh.paragraphs[0]
        ph.alignment = PP_ALIGN.CENTER
        rh = ph.add_run()
        rh.text = header
        rh.font.size = Pt(14)
        rh.font.bold = True
        rh.font.color.rgb = hex_to_rgb("#07090F" if highlight else accent)

        # Bullets
        tb = slide.shapes.add_textbox(x + Inches(0.1), col_top + Inches(0.7),
                                      col_w - Inches(0.2), H - col_top - Inches(1.3))
        tfb = tb.text_frame
        tfb.word_wrap = True
        for i, b in enumerate(bullets[:5]):
            pb = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
            rb = pb.add_run()
            rb.text = f"• {b}"
            rb.font.size = Pt(15)
            rb.font.color.rgb = hex_to_rgb(text)
            pb.space_after = Pt(6)

    add_column(left_x, left_header, left_bullets, highlight=False)
    add_column(right_x, right_header, right_bullets, highlight=highlight_right)

    # Vertical divider
    div = slide.shapes.add_connector(1, W / 2, col_top, W / 2, H - Inches(0.4))
    div.line.color.rgb = hex_to_rgb(accent)
    div.line.width = Emu(6350)

    return slide
